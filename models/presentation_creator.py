import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
from models.openai_langchain import OpenAiLC
import logging
import tempfile
import shutil
import re
from pptx.dml.color import RGBColor

'''
Mi sistema genera fotos con imagenes
libre de Derechos de Autor
'''

class PresentationCreator:
    def __init__(self):
        self.llmlc = OpenAiLC()
        logging.basicConfig(level=logging.DEBUG)

    def crear_presentacion(self, dataslide, tema, cantidad=1):
        if dataslide is None:
            print("Error: 'dataslide' no puede ser None")
            return None

        #print(f"Data slide viene de tools: {dataslide}")

        prs = Presentation()

        # Set the slide width and height to the default widescreen dimensions
        prs.slide_width = Inches(10.5)
        prs.slide_height = Inches(8.5)

        ahora = datetime.now()
        fecha_hora = ahora.strftime("fecha-%d-%m-%Y_hora-%H-%M-%S")

        for slide_info in dataslide:
            
            #print(f"slide_info: {slide_info}")

            slide_layout = prs.slide_layouts[5]  # Título Sólo layout
            slide = prs.slides.add_slide(slide_layout)

            # Agregar título
            title_box = slide.shapes.title
            title_box.text = slide_info["title"]

            # Ajustar el tamaño de la fuente del título
            title_text_frame = title_box.text_frame
            title_paragraph = title_text_frame.paragraphs[0]
            title_run = title_paragraph.runs[0]
            title_run.font.size = Pt(28)  # Cambia '24' al tamaño de fuente que desees

            # Cambiar el color del texto del título a azul
            title_run.font.color.rgb = RGBColor(0, 0, 255)  # RGB para azul

            # Agregar texto a la izquierda
            left_text = Inches(1)
            top_text = Inches(2)  # Ajusta según sea necesario
            width_text = Inches(5)
            height_text = Inches(6)
            txBox = slide.shapes.add_textbox(left_text, top_text, width_text, height_text)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.font.size = Pt(16) #ajusa el tamaño de la letra
            p.text = slide_info["content"]

            # Agregar imagen a la derecha
            try:
                image_url = self.llmlc.generar_imagen(slide_info["title"])
                response = requests.get(image_url, stream=True)
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file: # Crea un archivo temporal
                    shutil.copyfileobj(response.raw, temp_file) # Copia la imagen en el archivo temporal
                left_img = Inches(6)  # Ajusta según sea necesario
                top_img = Inches(2)  # Ajusta según sea necesario
                width_img = Inches(3)
                height_img = Inches(3)
                pic = slide.shapes.add_picture(temp_file.name, left_img, top_img, width_img, height_img)
                os.unlink(temp_file.name)  # Eliminar el archivo temporal
            except Exception as e:
                logging.error(f"Error al buscar imagen en OpenAi: {e}")
                print(f"Error al buscar imagen en OpenAi: {e}")
                image_url = None

        # Prepara el tema y la fecha_hora para ser seguros como nombres de archivo
        tema_safe = re.sub(r'[<>:"/\\|?*]', '', tema)  # Elimina caracteres no permitidos
        fecha_hora_safe = re.sub(r'[<>:"/\\|?*]', '', fecha_hora)  # Elimina caracteres no permitidos
        nombre_archivo = f"{tema_safe}-{fecha_hora_safe}.pptx"
        #nombre_archivo = f"{tema}-{fecha_hora}.pptx"
        ruta_guardado = os.path.join("RESULTADO", nombre_archivo)
        try:
            prs.save(ruta_guardado)
            if os.path.exists(ruta_guardado):
                return nombre_archivo
        except PermissionError:
            print("No se pudo crear el archivo, probablemente esté abierto.")
            return None
