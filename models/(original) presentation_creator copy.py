import os
import requests
from pptx import Presentation 
from pptx.util import Inches # Para cambiar el tamaño de la imagen
from pptx.util import Pt # Para cambiar el tamaño de la fuente
#from pptx.util import RGBColor # Para cambiar el color de la fuente
from datetime import datetime
from models.openai_langchain import OpenAiLC
import logging
import tempfile # Para crear directorio temporal
import shutil # Para eliminar directorio temporal

class PresentationCreator:
    def __init__(self):
        self.llmlc = OpenAiLC() 
        logging.basicConfig(level=logging.DEBUG)
        
    def crear_presentacion(self, dataslide, tema):
        if dataslide is None:
            print("Error: 'dataslide' no puede ser None")
            return None
        
        prs = Presentation()

        # Set the slide width and height to the default widescreen dimensions
        prs.slide_width = Inches(10.5)
        prs.slide_height = Inches(8.5)

        ahora = datetime.now()
        fecha_hora = ahora.strftime("fecha-%d-%m-%Y_hora-%H-%M-%S")
        
        for slide_info in dataslide:
            slide_layout = prs.slide_layouts[1]  #antes estaba el 5
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title and text
            title_box = slide.shapes.title
            title_box.text = slide_info["title"]
            left = Inches(1)
            top = Inches(1.5)
            width = height = Inches(8)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.font.size = Pt(14) 
            p.text = slide_info["content"]

            try:
                image_url = self.llmlc.generar_imagen(slide_info["title"])
                response = requests.get(image_url, stream=True)
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                    shutil.copyfileobj(response.raw, temp_file)
                left = Inches(1)
                top = Inches(2.5)
                pic = slide.shapes.add_picture(temp_file.name, left, top, width=Inches(3), height=Inches(3))
                os.unlink(temp_file.name)  # Eliminar el archivo temporal
            except Exception as e:
                logging.error(f"Error al buscar imagen en OpenAi: {e}")
                print(f"Error al buscar imagen en OpenAi: {e}")
                image_url = None

        nombre_archivo = f"{tema}-{fecha_hora}.pptx"
        ruta_guardado = os.path.join("RESULTADO", nombre_archivo)
        try:
            prs.save(ruta_guardado)
            if os.path.exists(ruta_guardado):
                return nombre_archivo
        except PermissionError:
            print("No se pudo crear el archivo, probablemente esté abierto.")
            return None

